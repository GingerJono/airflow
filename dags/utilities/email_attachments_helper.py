import base64
import logging
from io import BytesIO, StringIO
from json import dumps as json_serialise
from typing import Callable, Optional

from docx2txt import process as process_docx
from pandas import read_excel
from pptx import Presentation
from pymupdf import open as pdf_open

logger = logging.getLogger(__name__)


def _extract_text_pdf(pdf_file: BytesIO) -> str:
    with pdf_open(stream=pdf_file, filetype="pdf") as doc:
        pages = []
        for page_num in range(len(doc)):
            page = doc[page_num]
            pages.append({"page_number": page_num + 1, "text": page.get_text()})

    return json_serialise({"pages": pages})


def _extract_text_excel(excel_file: BytesIO) -> str:
    all_sheets = read_excel(excel_file, sheet_name=None)

    csv_outputs = {}
    for sheet_name, df in all_sheets.items():
        csv_buffer = StringIO()
        df.to_csv(csv_buffer, index=False)
        csv_outputs[sheet_name] = csv_buffer.getvalue()

    return json_serialise(csv_outputs)


def _extract_text_powerpoint(powerpoint_file: BytesIO) -> str:
    presentation = Presentation(powerpoint_file)

    all_slides = []
    for i, slide in enumerate(presentation.slides):
        shapes_on_slide = []
        for shape in slide.shapes:
            shapes_on_slide.append({"shape_text": shape.text})
        all_slides.append({"slide_number": i + 1, "shapes": shapes_on_slide})

    return json_serialise({"presentation": all_slides})


SUPPORTED_FILE_TYPES: dict[str, Callable[[BytesIO], str]] = {
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": process_docx,
    "application/pdf": _extract_text_pdf,
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": _extract_text_excel,
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": _extract_text_powerpoint,
}


def get_attachments_text(msgraph_attachments_response: dict) -> list[str]:
    attachments = []
    for msgraph_attachment in msgraph_attachments_response["value"]:
        attachment = _parse_attachment(msgraph_attachment)
        if attachment:
            attachments.append(json_serialise(attachment))

    return attachments


def _parse_attachment(msgraph_attachment: dict) -> Optional[dict]:
    file_type = msgraph_attachment["contentType"]

    if msgraph_attachment["isInline"]:
        logger.debug(
            "Attached file '%s' is an inline attachment and is therefore not included in the prompt",
            msgraph_attachment["name"],
        )
        return None

    if file_type not in SUPPORTED_FILE_TYPES:
        logger.warning(
            "Attached file '%s' is of an unsupported file type: %s",
            msgraph_attachment["name"],
            file_type,
        )
        return None

    attachment_bytes = base64.b64decode(msgraph_attachment["contentBytes"])

    attachment_file = BytesIO(attachment_bytes)

    attachment_text = SUPPORTED_FILE_TYPES[file_type](attachment_file)

    return {
        "original_file_type": file_type,
        "file_name": msgraph_attachment["name"],
        "content": attachment_text,
    }
